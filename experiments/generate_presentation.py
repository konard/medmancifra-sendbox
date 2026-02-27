"""
Generate PPTX presentation for AI-Platform for Automating Bioequivalence Research (Oncology).
Based on issue #5: https://github.com/medmancifra/sendbox/issues/5
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Color palette
BLUE_DARK = RGBColor(0x1A, 0x23, 0x7E)   # Dark blue
BLUE_MED = RGBColor(0x15, 0x65, 0xC0)    # Medium blue
BLUE_LIGHT = RGBColor(0x42, 0xA5, 0xF5)  # Light blue
PURPLE = RGBColor(0x6A, 0x1B, 0x9A)      # Purple
GREEN = RGBColor(0x2E, 0x7D, 0x32)       # Green
TEAL = RGBColor(0x00, 0x89, 0x7B)        # Teal
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x21, 0x21, 0x21)
GRAY_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
GRAY_MED = RGBColor(0xE0, 0xE0, 0xE0)
ACCENT_ORANGE = RGBColor(0xEF, 0x6C, 0x00)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)


def set_background(slide, color: RGBColor):
    """Set slide background to a solid color."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_background(slide, color1: RGBColor, color2: RGBColor):
    """Set slide background to gradient (approximated with solid for compatibility)."""
    # python-pptx doesn't directly support gradient backgrounds well,
    # so we approximate with a solid color (midpoint blend)
    r = (color1.rgb >> 16 & 0xFF + color2.rgb >> 16 & 0xFF) // 2
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color1


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_colored_box(slide, left, top, width, height, fill_color: RGBColor,
                    text=None, text_color=WHITE, font_size=14, bold=False,
                    align=PP_ALIGN.CENTER):
    """Add a colored rectangle shape with optional text."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color

    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color

    return shape


def add_slide_number(slide, number, total, color=WHITE):
    """Add slide number to bottom right."""
    add_text_box(
        slide, f"{number} / {total}",
        left=Inches(8.8), top=Inches(6.8),
        width=Inches(1.0), height=Inches(0.3),
        font_size=10, color=color,
        align=PP_ALIGN.RIGHT
    )


SLIDE_COUNT = 9  # Total slides


def create_title_slide(prs):
    """Slide 1: Title / Main slide"""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Gradient background (blue to purple - approximate with dark blue)
    set_background(slide, BLUE_DARK)

    # Top accent bar
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(0.08),
                    fill_color=BLUE_LIGHT)

    # Main title
    add_text_box(
        slide,
        "AI-Platform for Automating\nBioequivalence Research",
        left=Inches(0.5), top=Inches(1.0),
        width=Inches(9.0), height=Inches(1.8),
        font_size=40, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Subtitle tag: Oncology
    add_colored_box(
        slide,
        left=Inches(3.5), top=Inches(2.9),
        width=Inches(3.0), height=Inches(0.5),
        fill_color=PURPLE,
        text="ONCOLOGY", text_color=WHITE, font_size=16, bold=True,
        align=PP_ALIGN.CENTER
    )

    # Subtitle
    add_text_box(
        slide,
        "Digital Health  |  AI-Architecture  |  Hackathon / Invest-Pitch",
        left=Inches(0.5), top=Inches(3.6),
        width=Inches(9.0), height=Inches(0.6),
        font_size=16, color=BLUE_LIGHT,
        align=PP_ALIGN.CENTER
    )

    # Slogan
    add_text_box(
        slide,
        '"We make doctors\' lives more cloud-based and IT-friendly"',
        left=Inches(0.5), top=Inches(4.4),
        width=Inches(9.0), height=Inches(0.6),
        font_size=14, italic=True, color=GRAY_MED,
        align=PP_ALIGN.CENTER
    )

    # Icon placeholders (text symbols)
    icons_text = "🚀  ☁️  🤖  💊  📊"
    add_text_box(
        slide,
        icons_text,
        left=Inches(2.5), top=Inches(5.2),
        width=Inches(5.0), height=Inches(0.6),
        font_size=24, color=WHITE,
        align=PP_ALIGN.CENTER
    )

    # Bottom bar
    add_colored_box(slide, Inches(0), Inches(6.8), Inches(10), Inches(0.7),
                    fill_color=PURPLE,
                    text="konard / medmancifra  ·  medmancifra/sendbox  ·  2024",
                    text_color=GRAY_MED, font_size=11)

    add_slide_number(slide, 1, SLIDE_COUNT, WHITE)
    return slide


def create_problem_slide(prs):
    """Slide 2: Problem"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, WHITE)

    # Header bar
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.2),
                    fill_color=BLUE_DARK,
                    text="Bioequivalence Problem in Oncology",
                    text_color=WHITE, font_size=26, bold=True,
                    align=PP_ALIGN.CENTER)

    # Problem items with icons
    problems = [
        ("📈", "High intra-subject variability", "CV >30% — requires RSABE approach instead of standard ABE"),
        ("⏱️", "Long half-life drugs", "t½ > 24–72h — complex washout planning, accumulation risk"),
        ("📋", "Multiple regulatory frameworks", "EMA/FDA/EEC No.85 require different criteria — manual reconciliation"),
        ("🍽️", "Fasting/Fed requirements", "Many onco drugs require two separate studies for food effects"),
        ("👤", "Human factor in planning", "Manual calculations → errors, inconsistencies, audit failures"),
    ]

    y = Inches(1.4)
    for icon, title, desc in problems:
        # Icon circle
        add_colored_box(slide,
                        left=Inches(0.3), top=y,
                        width=Inches(0.7), height=Inches(0.7),
                        fill_color=BLUE_LIGHT,
                        text=icon, font_size=20, align=PP_ALIGN.CENTER)
        # Title
        add_text_box(slide, title,
                     left=Inches(1.2), top=y,
                     width=Inches(3.5), height=Inches(0.4),
                     font_size=14, bold=True, color=BLUE_DARK)
        # Description
        add_text_box(slide, desc,
                     left=Inches(1.2), top=y + Inches(0.35),
                     width=Inches(8.4), height=Inches(0.35),
                     font_size=11, color=BLACK)
        y += Inches(0.85)

    # Bottom "Result" box
    add_colored_box(slide,
                    left=Inches(0.3), top=Inches(5.7),
                    width=Inches(9.4), height=Inches(0.7),
                    fill_color=ACCENT_RED,
                    text="Result: 2–5 days of manual planning per study → errors, regulatory risks, delays",
                    text_color=WHITE, font_size=13, bold=True)

    add_slide_number(slide, 2, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_solution_slide(prs):
    """Slide 3: Solution"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, GRAY_LIGHT)

    # Header bar
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.2),
                    fill_color=BLUE_MED,
                    text="Regulatory-Aware AI Planning Engine",
                    text_color=WHITE, font_size=26, bold=True,
                    align=PP_ALIGN.CENTER)

    # Four solution blocks
    blocks = [
        (BLUE_DARK, "🎯 Deterministic Core",
         "No LLM in math\nOwen's method\nISO/GOST compliant\ncalculations"),
        (PURPLE, "🤖 AI Explanation Layer",
         "GigaChat RAG\nNatural language\nexplanations for\neach decision"),
        (TEAL, "🏥 Onco-Specific Logic",
         "RSABE module\nCYP interactions\nHigh-variability\nprotocols"),
        (GREEN, "📄 Synopsis Generation",
         "Auto-generated\nMarkdown/JSON\nRegulatory-ready\nstructure"),
    ]

    x_positions = [Inches(0.3), Inches(2.8), Inches(5.3), Inches(7.8)]
    for i, (color, title, desc) in enumerate(blocks):
        x = x_positions[i]
        # Block header
        add_colored_box(slide,
                        left=x, top=Inches(1.4),
                        width=Inches(2.3), height=Inches(0.8),
                        fill_color=color,
                        text=title, font_size=12, bold=True, align=PP_ALIGN.CENTER)
        # Block content
        add_colored_box(slide,
                        left=x, top=Inches(2.2),
                        width=Inches(2.3), height=Inches(2.0),
                        fill_color=WHITE,
                        text=desc, text_color=BLACK, font_size=11,
                        align=PP_ALIGN.LEFT)

    # Arrow flow: Input → Platform → Output
    add_text_box(slide, "📥 INPUT\n(CV, t½,\nfood effect)",
                 left=Inches(0.2), top=Inches(4.5),
                 width=Inches(2.0), height=Inches(1.2),
                 font_size=12, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    add_text_box(slide, "➡️➡️➡️",
                 left=Inches(2.3), top=Inches(4.8),
                 width=Inches(1.0), height=Inches(0.6),
                 font_size=20, color=BLUE_MED, align=PP_ALIGN.CENTER)

    add_colored_box(slide,
                    left=Inches(3.3), top=Inches(4.4),
                    width=Inches(3.4), height=Inches(1.2),
                    fill_color=BLUE_MED,
                    text="AI-Platform\n(30-90 min automated)",
                    text_color=WHITE, font_size=13, bold=True)

    add_text_box(slide, "➡️➡️➡️",
                 left=Inches(6.8), top=Inches(4.8),
                 width=Inches(1.0), height=Inches(0.6),
                 font_size=20, color=BLUE_MED, align=PP_ALIGN.CENTER)

    add_text_box(slide, "📤 OUTPUT\n(Design + N +\nSynopsis)",
                 left=Inches(7.8), top=Inches(4.5),
                 width=Inches(2.0), height=Inches(1.2),
                 font_size=12, color=GREEN, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 3, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_architecture_slide(prs):
    """Slide 4: System Architecture"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, WHITE)

    # Header
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.0),
                    fill_color=BLUE_DARK,
                    text="Platform Architecture",
                    text_color=WHITE, font_size=26, bold=True,
                    align=PP_ALIGN.CENTER)

    # Frontend box
    add_colored_box(slide,
                    left=Inches(0.3), top=Inches(1.2),
                    width=Inches(3.5), height=Inches(4.5),
                    fill_color=RGBColor(0xE3, 0xF2, 0xFD),
                    text="", align=PP_ALIGN.CENTER)

    add_text_box(slide, "FRONTEND\n(React + TypeScript)",
                 left=Inches(0.4), top=Inches(1.3),
                 width=Inches(3.3), height=Inches(0.7),
                 font_size=14, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

    frontend_items = [
        "🖥️ Calculator Form",
        "📊 Results Panel",
        "💬 AI Chat (RAG)",
        "📄 Synopsis Export",
        "   (MD / JSON)",
    ]
    y = Inches(2.1)
    for item in frontend_items:
        add_text_box(slide, item,
                     left=Inches(0.5), top=y,
                     width=Inches(3.1), height=Inches(0.35),
                     font_size=12, color=BLUE_DARK)
        y += Inches(0.35)

    # Arrow
    add_text_box(slide, "⟺\nHTTP\nREST",
                 left=Inches(3.9), top=Inches(2.7),
                 width=Inches(0.8), height=Inches(1.2),
                 font_size=11, color=GRAY_MED, align=PP_ALIGN.CENTER)

    # Backend box
    add_colored_box(slide,
                    left=Inches(4.8), top=Inches(1.2),
                    width=Inches(4.9), height=Inches(4.5),
                    fill_color=RGBColor(0xE8, 0xF5, 0xE9),
                    text="", align=PP_ALIGN.CENTER)

    add_text_box(slide, "BACKEND\n(Kotlin + Quarkus 3.8)",
                 left=Inches(4.9), top=Inches(1.3),
                 width=Inches(4.7), height=Inches(0.7),
                 font_size=14, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    backend_modules = [
        (BLUE_MED, "Design+N Engine", "DesignEngine.kt · SampleSizeEngine.kt"),
        (PURPLE, "RSABE Module", "BECalculationService.kt (Owen's method)"),
        (TEAL, "PK Intelligence", "washout, accumulation, t½ logic"),
        (ACCENT_ORANGE, "Regulatory RAG", "GigaChat + EEC No.85 / EMA / FDA"),
    ]

    y = Inches(2.1)
    for color, title, detail in backend_modules:
        add_colored_box(slide,
                        left=Inches(4.9), top=y,
                        width=Inches(4.6), height=Inches(0.6),
                        fill_color=color,
                        text=f"{title}  —  {detail}",
                        text_color=WHITE, font_size=10, bold=False,
                        align=PP_ALIGN.LEFT)
        y += Inches(0.72)

    # Decision flow
    add_text_box(slide, "Decision flow:",
                 left=Inches(0.3), top=Inches(5.9),
                 width=Inches(1.5), height=Inches(0.35),
                 font_size=10, bold=True, color=BLUE_DARK)
    add_text_box(slide,
                 "CV_intra ≤30% → 2×2  |  >30% → RSABE → replicate/parallel → washout → fasting/fed → N → Synopsis",
                 left=Inches(1.9), top=Inches(5.9),
                 width=Inches(7.8), height=Inches(0.35),
                 font_size=10, color=BLACK)

    # Database
    add_colored_box(slide,
                    left=Inches(4.8), top=Inches(5.7),
                    width=Inches(2.2), height=Inches(0.55),
                    fill_color=RGBColor(0xFF, 0xF3, 0xE0),
                    text="🗄️ PostgreSQL",
                    text_color=ACCENT_ORANGE, font_size=11, bold=True)

    add_slide_number(slide, 4, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_usecase_slide(prs):
    """Slide 5: Use Case"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, GRAY_LIGHT)

    # Header
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.0),
                    fill_color=TEAL,
                    text="Example: High-Variability Oncology Drug",
                    text_color=WHITE, font_size=24, bold=True,
                    align=PP_ALIGN.CENTER)

    # Input parameters box
    add_colored_box(slide,
                    left=Inches(0.3), top=Inches(1.2),
                    width=Inches(4.2), height=Inches(2.5),
                    fill_color=WHITE,
                    text="", align=PP_ALIGN.CENTER)
    add_text_box(slide, "📥 INPUT PARAMETERS",
                 left=Inches(0.4), top=Inches(1.3),
                 width=Inches(4.0), height=Inches(0.45),
                 font_size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    params = [
        ("CV_intra", "42%", "→ High variability (RSABE)"),
        ("t½", "48 h", "→ Long washout required"),
        ("Food Effect", "Confirmed", "→ 2 studies: fasting + fed"),
    ]
    y = Inches(1.85)
    for param, value, note in params:
        add_text_box(slide, f"• {param}: ", left=Inches(0.5), top=y,
                     width=Inches(1.2), height=Inches(0.4), font_size=12, color=BLACK)
        add_text_box(slide, value, left=Inches(1.7), top=y,
                     width=Inches(0.8), height=Inches(0.4),
                     font_size=13, bold=True, color=ACCENT_RED)
        add_text_box(slide, note, left=Inches(2.5), top=y,
                     width=Inches(2.0), height=Inches(0.4),
                     font_size=10, italic=True, color=GRAY_MED)
        y += Inches(0.45)

    # Arrow
    add_text_box(slide, "➡️",
                 left=Inches(4.6), top=Inches(2.3),
                 width=Inches(0.8), height=Inches(0.6),
                 font_size=28, align=PP_ALIGN.CENTER)

    # Platform output box
    add_colored_box(slide,
                    left=Inches(5.5), top=Inches(1.2),
                    width=Inches(4.2), height=Inches(2.5),
                    fill_color=RGBColor(0xE8, 0xF5, 0xE9),
                    text="", align=PP_ALIGN.CENTER)
    add_text_box(slide, "📤 PLATFORM RESULTS",
                 left=Inches(5.6), top=Inches(1.3),
                 width=Inches(4.0), height=Inches(0.45),
                 font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    results = [
        "✅  Study Design: 4-period replicate",
        "✅  Studies: 2 (fasting + fed)",
        "✅  N randomized: 64",
        "✅  N screening: 92",
        "✅  Synopsis: auto-generated",
    ]
    y = Inches(1.85)
    for result in results:
        add_text_box(slide, result,
                     left=Inches(5.7), top=y,
                     width=Inches(3.8), height=Inches(0.4),
                     font_size=12, color=BLACK)
        y += Inches(0.42)

    # Time reduction bar
    add_text_box(slide, "⏱️ Time Reduction",
                 left=Inches(0.3), top=Inches(3.9),
                 width=Inches(3.0), height=Inches(0.4),
                 font_size=13, bold=True, color=BLUE_DARK)

    # Manual bar
    add_text_box(slide, "Manual planning: 2–5 days",
                 left=Inches(0.3), top=Inches(4.4),
                 width=Inches(3.0), height=Inches(0.35),
                 font_size=11, color=ACCENT_RED)
    add_colored_box(slide, Inches(3.4), Inches(4.4), Inches(6.2), Inches(0.35),
                    fill_color=ACCENT_RED, text="", align=PP_ALIGN.CENTER)

    # AI bar
    add_text_box(slide, "AI-Platform: 30–90 min",
                 left=Inches(0.3), top=Inches(4.9),
                 width=Inches(3.0), height=Inches(0.35),
                 font_size=11, color=GREEN)
    add_colored_box(slide, Inches(3.4), Inches(4.9), Inches(2.2), Inches(0.35),
                    fill_color=GREEN, text="", align=PP_ALIGN.CENTER)

    # Saving badge
    add_colored_box(slide,
                    left=Inches(0.3), top=Inches(5.5),
                    width=Inches(9.4), height=Inches(0.7),
                    fill_color=TEAL,
                    text="⚡  60–70% time savings  ·  Eliminate manual errors  ·  Regulatory-ready output in minutes",
                    text_color=WHITE, font_size=13, bold=True)

    add_slide_number(slide, 5, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_competitive_slide(prs):
    """Slide 6: Competitive Advantages"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, WHITE)

    # Header
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.0),
                    fill_color=BLUE_DARK,
                    text="Competitive Advantages",
                    text_color=WHITE, font_size=26, bold=True,
                    align=PP_ALIGN.CENTER)

    # Table headers
    headers = ["Solution Type", "Limitation", "Our Advantage"]
    header_colors = [BLUE_MED, ACCENT_RED, GREEN]
    widths = [Inches(2.8), Inches(3.2), Inches(3.6)]
    x_positions = [Inches(0.3), Inches(3.1), Inches(6.3)]

    y = Inches(1.1)
    for i, (header, color, w, x) in enumerate(zip(headers, header_colors, widths, x_positions)):
        add_colored_box(slide, x, y, w, Inches(0.5),
                        fill_color=color, text=header,
                        text_color=WHITE, font_size=13, bold=True)

    # Table rows
    rows = [
        ("Manual calculation\n(Excel/SAS)", "Human errors\nNot regulatory-aware\n2–5 days per study",
         "✅ Automated engine\n✅ EMA/FDA compliant\n✅ 30–90 min"),
        ("Generic BE software\n(WinNonlin, Phoenix)", "No oncology logic\nNo regulatory AI\nExpensive licenses",
         "✅ Onco-specific RSABE\n✅ Built-in RAG\n✅ Cloud SaaS model"),
        ("CRO consulting\n(manual expert work)", "High cost per study\nNon-reproducible\nSlow turnaround",
         "✅ Deterministic core\n✅ Reproducible results\n✅ Audit trail"),
    ]

    y = Inches(1.65)
    row_colors = [RGBColor(0xF8, 0xF9, 0xFA), RGBColor(0xF1, 0xF3, 0xF4), RGBColor(0xE8, 0xEA, 0xED)]
    for row_idx, (sol, lim, adv) in enumerate(rows):
        row_h = Inches(1.1)
        bg = row_colors[row_idx]
        for i, (text, w, x) in enumerate(zip([sol, lim, adv], widths, x_positions)):
            add_colored_box(slide, x, y, w, row_h,
                            fill_color=bg, text=text,
                            text_color=BLACK, font_size=10,
                            align=PP_ALIGN.LEFT)
        y += row_h + Inches(0.05)

    # Bottom summary
    add_colored_box(slide,
                    left=Inches(0.3), top=Inches(5.6),
                    width=Inches(9.4), height=Inches(0.7),
                    fill_color=BLUE_DARK,
                    text="We combine statistics + regulatory expertise + AI explanations in one platform",
                    text_color=WHITE, font_size=14, bold=True)

    add_slide_number(slide, 6, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_investment_slide(prs):
    """Slide 7: Investment Logic"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, BLUE_DARK)

    # Top accent
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(0.08),
                    fill_color=GREEN)

    # Header
    add_text_box(slide, "Investment Logic",
                 left=Inches(0.5), top=Inches(0.2),
                 width=Inches(9.0), height=Inches(0.9),
                 font_size=30, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Three investment blocks
    blocks = [
        (BLUE_MED, "📊 BE is Standardized",
         "→ Algorithmic\n\nHighly regulated domain\nwith deterministic rules\n→ perfect for automation"),
        (PURPLE, "🏥 Oncology is Complex",
         "→ High Automation Value\n\nHigh-variability drugs\nrequire expert-level logic\n→ AI provides leverage"),
        (GREEN, "🚀 AI Infrastructure",
         "Not just a calculator\n\nFull regulatory stack:\nRAG + RSABE + PK +\nSynopsis generation"),
    ]

    x_positions = [Inches(0.3), Inches(3.5), Inches(6.7)]
    for i, (color, title, desc) in enumerate(blocks):
        x = x_positions[i]
        add_colored_box(slide, x, Inches(1.3), Inches(3.0), Inches(0.7),
                        fill_color=color, text=title,
                        text_color=WHITE, font_size=13, bold=True)
        add_colored_box(slide, x, Inches(2.0), Inches(3.0), Inches(2.0),
                        fill_color=RGBColor(0x1E, 0x2A, 0x40),
                        text=desc, text_color=WHITE, font_size=11,
                        align=PP_ALIGN.LEFT)

    # ROI section
    add_text_box(slide, "💰 ROI / Time Savings",
                 left=Inches(0.3), top=Inches(4.2),
                 width=Inches(3.0), height=Inches(0.45),
                 font_size=16, bold=True, color=GREEN)

    roi_data = [
        ("Savings per study", "~88,000 RUB"),
        ("30 studies/year", "~2.6M RUB saved"),
        ("Implementation cost", "~3M RUB"),
        ("Breakeven", "~12 months"),
        ("5-year ROI", "400–600%"),
    ]

    y = Inches(4.75)
    for label, value in roi_data:
        add_text_box(slide, f"{label}:", left=Inches(0.3), top=y,
                     width=Inches(2.5), height=Inches(0.35),
                     font_size=11, color=GRAY_MED)
        add_text_box(slide, value, left=Inches(2.9), top=y,
                     width=Inches(2.0), height=Inches(0.35),
                     font_size=12, bold=True, color=GREEN)
        y += Inches(0.35)

    # Upward arrow / growth symbol
    add_text_box(slide, "📈",
                 left=Inches(5.5), top=Inches(4.2),
                 width=Inches(4.0), height=Inches(2.3),
                 font_size=72, color=GREEN, align=PP_ALIGN.CENTER)

    add_slide_number(slide, 7, SLIDE_COUNT, WHITE)
    return slide


def create_economics_slide(prs):
    """Slide 8: Economics — Russian BE Market"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, WHITE)

    # Header
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(1.0),
                    fill_color=TEAL,
                    text="Russian BE Market — Economics",
                    text_color=WHITE, font_size=24, bold=True,
                    align=PP_ALIGN.CENTER)

    # Market table
    companies = [
        ("Synergy Research Group", "900–1,200M"),
        ("Medical Development Agency", "290–380M"),
        ("iPharma", "360–480M"),
        ("CPD CRO", "180–240M"),
        ("QBio Research", "96–128M"),
        ("OCT CRO", "165–220M"),
        ("R&D Pharma", "72–96M"),
    ]

    # Table header
    add_colored_box(slide, Inches(0.3), Inches(1.1), Inches(5.0), Inches(0.45),
                    fill_color=TEAL, text="CRO Company",
                    text_color=WHITE, font_size=12, bold=True)
    add_colored_box(slide, Inches(5.4), Inches(1.1), Inches(4.2), Inches(0.45),
                    fill_color=TEAL, text="Revenue/Year (RUB)",
                    text_color=WHITE, font_size=12, bold=True)

    row_colors_alt = [RGBColor(0xF0, 0xFD, 0xFA), RGBColor(0xE0, 0xF7, 0xF4)]
    y = Inches(1.6)
    for i, (company, revenue) in enumerate(companies):
        bg = row_colors_alt[i % 2]
        add_colored_box(slide, Inches(0.3), y, Inches(5.0), Inches(0.42),
                        fill_color=bg, text=company,
                        text_color=BLACK, font_size=11, align=PP_ALIGN.LEFT)
        add_colored_box(slide, Inches(5.4), y, Inches(4.2), Inches(0.42),
                        fill_color=bg, text=revenue,
                        text_color=TEAL, font_size=12, bold=True,
                        align=PP_ALIGN.CENTER)
        y += Inches(0.45)

    # Cost comparison
    add_text_box(slide, "Cost per BE Study",
                 left=Inches(0.3), top=Inches(5.0),
                 width=Inches(4.0), height=Inches(0.4),
                 font_size=14, bold=True, color=BLUE_DARK)

    comparisons = [
        ("Manual planning (expert CRO):", "96,000 RUB", ACCENT_RED),
        ("AI-Platform (automated):", "8,000 RUB", GREEN),
        ("Savings per study:", "88,000 RUB", TEAL),
    ]
    y = Inches(5.5)
    for label, value, color in comparisons:
        add_text_box(slide, label, left=Inches(0.3), top=y,
                     width=Inches(3.5), height=Inches(0.35),
                     font_size=11, color=BLACK)
        add_text_box(slide, value, left=Inches(3.9), top=y,
                     width=Inches(2.0), height=Inches(0.35),
                     font_size=12, bold=True, color=color)
        y += Inches(0.38)

    add_slide_number(slide, 8, SLIDE_COUNT, BLUE_DARK)
    return slide


def create_team_slide(prs):
    """Slide 9: Team"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    set_background(slide, BLUE_DARK)

    # Top accent
    add_colored_box(slide, Inches(0), Inches(0), Inches(10), Inches(0.08),
                    fill_color=ACCENT_ORANGE)

    # Header
    add_text_box(slide, "Team",
                 left=Inches(0.5), top=Inches(0.2),
                 width=Inches(9.0), height=Inches(0.9),
                 font_size=32, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Team members
    members = [
        ("🏥", "Clinical Pharmacologist", "BE study design, PK analysis,\nregulatory submissions"),
        ("💻", "Backend Engineer", "Kotlin/Quarkus, calculation\nengines, REST APIs"),
        ("🎨", "Frontend Developer", "React/TypeScript, UX/UI,\ncalculator interface"),
        ("🤖", "AI/ML Engineer", "GigaChat RAG, NLP,\nexplanation layer"),
    ]

    x_positions = [Inches(0.3), Inches(2.7), Inches(5.1), Inches(7.5)]

    for i, (icon, role, skills) in enumerate(members):
        x = x_positions[i]
        # Avatar circle
        add_colored_box(slide, x + Inches(0.4), Inches(1.3),
                        Inches(1.4), Inches(1.4),
                        fill_color=BLUE_MED,
                        text=icon, font_size=28, align=PP_ALIGN.CENTER)
        # Role title
        add_text_box(slide, role,
                     left=x, top=Inches(2.85),
                     width=Inches(2.2), height=Inches(0.6),
                     font_size=12, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        # Skills
        add_text_box(slide, skills,
                     left=x, top=Inches(3.5),
                     width=Inches(2.2), height=Inches(0.8),
                     font_size=10, color=GRAY_MED,
                     align=PP_ALIGN.CENTER)

    # GitHub repos section
    add_text_box(slide, "🔗 Project Links",
                 left=Inches(0.3), top=Inches(4.5),
                 width=Inches(3.0), height=Inches(0.4),
                 font_size=14, bold=True, color=ACCENT_ORANGE)

    links = [
        ("Repository", "github.com/medmancifra/sendbox"),
        ("Issue", "github.com/medmancifra/sendbox/issues/5"),
        ("PR / Demo", "github.com/medmancifra/sendbox/pull/6"),
    ]
    y = Inches(5.0)
    for label, url in links:
        add_text_box(slide, f"{label}: {url}",
                     left=Inches(0.3), top=y,
                     width=Inches(7.0), height=Inches(0.35),
                     font_size=11, color=GRAY_MED)
        y += Inches(0.35)

    # Bottom bar
    add_colored_box(slide, Inches(0), Inches(6.8), Inches(10), Inches(0.7),
                    fill_color=PURPLE,
                    text="AI-Platform for BE Research (Oncology)  ·  2024  ·  Made with 💊 + 🤖",
                    text_color=GRAY_MED, font_size=11)

    add_slide_number(slide, 9, SLIDE_COUNT, WHITE)
    return slide


def main():
    output_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                              "docs", "presentation")
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation()

    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")
    create_title_slide(prs)
    print("  ✅ Slide 1: Title")
    create_problem_slide(prs)
    print("  ✅ Slide 2: Problem")
    create_solution_slide(prs)
    print("  ✅ Slide 3: Solution")
    create_architecture_slide(prs)
    print("  ✅ Slide 4: Architecture")
    create_usecase_slide(prs)
    print("  ✅ Slide 5: Use Case")
    create_competitive_slide(prs)
    print("  ✅ Slide 6: Competitive Advantages")
    create_investment_slide(prs)
    print("  ✅ Slide 7: Investment Logic")
    create_economics_slide(prs)
    print("  ✅ Slide 8: Economics")
    create_team_slide(prs)
    print("  ✅ Slide 9: Team")

    pptx_path = os.path.join(output_dir, "BE_AI_Platform_Presentation.pptx")
    prs.save(pptx_path)
    print(f"\n✅ Presentation saved: {pptx_path}")
    print(f"   File size: {os.path.getsize(pptx_path) / 1024:.1f} KB")
    return pptx_path


if __name__ == "__main__":
    main()
